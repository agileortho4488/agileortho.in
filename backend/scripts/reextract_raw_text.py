"""
Re-extract raw page text for Batch 1 files (1-25).
Adds _raw_text_by_page, _tables_by_page, page_extractions to existing raw_extraction files.
Does NOT touch structured_drafts — those stay as-is.
"""
import json
import os
import pdfplumber
import pytesseract
from datetime import datetime, timezone

BASE_DIR = "/app/backend/brochure_intelligence"
RAW_DIR = os.path.join(BASE_DIR, "raw_extractions")
BROCHURE_DIR = "/tmp/zoho_brochures"
NOW = datetime.now(timezone.utc).isoformat()


def extract_pdf_pages(filepath, dpi=300):
    pages = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            method = "pdfplumber_direct"
            tables = page.extract_tables()
            table_data = tables if tables else []

            if not text or len(text.strip()) < 50:
                try:
                    img = page.to_image(resolution=dpi)
                    text = pytesseract.image_to_string(img.original)
                    method = f"pytesseract_OCR_{dpi}dpi"
                except Exception as e:
                    text = f"[OCR FAILED: {str(e)}]"
                    method = "ocr_failed"

            pages.append({
                "page_number": i + 1,
                "text": text.strip() if text else "",
                "text_length": len(text.strip()) if text else 0,
                "extraction_method": method,
                "has_tables": len(table_data) > 0,
                "tables": table_data,
                "is_blank": len(text.strip()) < 10 if text else True,
            })
    return pages


def extract_pptx_pages(filepath):
    from pptx import Presentation
    pages = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        text = "\n".join(texts)
        pages.append({
            "page_number": i + 1,
            "text": text,
            "text_length": len(text),
            "extraction_method": "python-pptx",
            "has_tables": False,
            "tables": [],
            "is_blank": len(text) < 10,
        })
    return pages


def process_file(raw_path):
    with open(raw_path) as f:
        data = json.load(f)

    eid = data.get("extraction_id", "")
    source_file = data.get("source_file", "")
    file_type = data.get("file_type", "pdf")

    if data.get("_raw_evidence_status") == "PRESENT":
        return False  # Already has raw text

    brochure_path = os.path.join(BROCHURE_DIR, source_file)
    if not os.path.exists(brochure_path):
        print(f"  {eid}: SKIPPED — source file not found: {source_file}")
        data["_raw_evidence_status"] = "SOURCE_NOT_FOUND"
        with open(raw_path, 'w') as f:
            json.dump(data, f, indent=2)
        return False

    # Extract raw pages
    if file_type == "pdf":
        pages = extract_pdf_pages(brochure_path)
    elif file_type in ("pptx", "ppt"):
        pages = extract_pptx_pages(brochure_path)
    elif file_type in ("docx", "xlsx"):
        # Handle corrupted or special formats
        print(f"  {eid}: SKIPPED — {file_type} format (needs special handling)")
        data["_raw_evidence_status"] = "NEEDS_SPECIAL_HANDLER"
        with open(raw_path, 'w') as f:
            json.dump(data, f, indent=2)
        return False
    else:
        print(f"  {eid}: SKIPPED — unsupported type: {file_type}")
        return False

    # Store raw evidence
    data["_raw_text_by_page"] = {
        str(p["page_number"]): p["text"] for p in pages
    }
    data["_tables_by_page"] = {
        str(p["page_number"]): p["tables"] for p in pages if p["tables"]
    }
    data["page_extractions"] = [
        {
            "source_file": source_file,
            "file_id": data.get("file_id", ""),
            "page_number": p["page_number"],
            "parser_used": p["extraction_method"],
            "extraction_method": p["extraction_method"],
            "text_length": p["text_length"],
            "is_blank": p["is_blank"],
            "has_tables": p["has_tables"]
        }
        for p in pages
    ]
    data["total_pages"] = len(pages)
    data["_raw_evidence_status"] = "PRESENT"
    data["_raw_extracted_at"] = NOW
    data["_has_raw_text"] = True
    data["_has_tables"] = any(p["has_tables"] for p in pages)
    data["_has_page_metadata"] = True

    total_text = sum(p["text_length"] for p in pages)
    blank_pages = sum(1 for p in pages if p["is_blank"])
    parsers = set(p["extraction_method"] for p in pages)

    with open(raw_path, 'w') as f:
        json.dump(data, f, indent=2)

    print(f"  {eid}: {source_file[:40]} — {len(pages)} pages, {total_text} chars, {blank_pages} blank, parsers: {parsers}")
    return True


if __name__ == "__main__":
    files = sorted([f for f in os.listdir(RAW_DIR) if f.endswith('.json')])
    processed = 0
    skipped = 0

    for fname in files:
        raw_path = os.path.join(RAW_DIR, fname)
        if process_file(raw_path):
            processed += 1
        else:
            skipped += 1

    print(f"\nRe-extraction complete: {processed} files updated, {skipped} skipped")
