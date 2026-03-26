"""
Batch Extraction Script — Files 26-50
Follows strict rules:
1. No file is complete until every page is checked
2. Every page stays married to its original brochure source
3. Actual SKU/spec data preserved, not just counts
"""
import json
import os
import hashlib
import pdfplumber
import pytesseract
from datetime import datetime, timezone
from pptx import Presentation

BASE_DIR = "/app/backend/brochure_intelligence"
BROCHURE_DIR = "/tmp/zoho_brochures"
RAW_DIR = os.path.join(BASE_DIR, "raw_extractions")
NOW = datetime.now(timezone.utc).isoformat()


def load_file_index():
    with open(os.path.join(BASE_DIR, "manifests/file_index.json")) as f:
        return json.load(f)


def extract_pdf_pages(filepath, dpi=300):
    """Extract text from every page of a PDF using direct text + OCR fallback"""
    pages = []
    with pdfplumber.open(filepath) as pdf:
        for i, page in enumerate(pdf.pages):
            # Step 1: Try direct text extraction
            text = page.extract_text()
            method = "pdfplumber_direct"

            # Step 2: Try table extraction
            tables = page.extract_tables()
            table_data = []
            if tables:
                for table in tables:
                    table_data.append(table)

            # Step 3: OCR fallback if direct text is insufficient
            if not text or len(text.strip()) < 50:
                try:
                    img = page.to_image(resolution=dpi)
                    text = pytesseract.image_to_string(img.original)
                    method = f"pytesseract_OCR_{dpi}dpi"
                except Exception as e:
                    text = f"[OCR FAILED: {str(e)}]"
                    method = "ocr_failed"

            pages.append({
                "page_number": i + 1,
                "text": text.strip() if text else "",
                "text_length": len(text.strip()) if text else 0,
                "extraction_method": method,
                "has_tables": len(table_data) > 0,
                "tables": table_data,
                "is_blank": len(text.strip()) < 10 if text else True,
            })

    return pages


def extract_pptx_pages(filepath):
    """Extract text from every slide of a PPTX"""
    pages = []
    prs = Presentation(filepath)
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        text = "\n".join(texts)
        pages.append({
            "page_number": i + 1,
            "text": text,
            "text_length": len(text),
            "extraction_method": "python-pptx",
            "has_tables": False,
            "tables": [],
            "is_blank": len(text) < 10,
        })
    return pages


def classify_content(pages, filename):
    """Classify the content type based on page text"""
    all_text = " ".join(p["text"] for p in pages).lower()

    if any(kw in all_text for kw in ["sku", "catalog", "ordering", "part number", "product code"]):
        return "product_catalog"
    if any(kw in all_text for kw in ["brochure", "specifications", "features", "technical"]):
        return "product_brochure"
    if any(kw in all_text for kw in ["questionnaire", "survey", "feedback"]):
        return "reference_doc"
    if any(kw in all_text for kw in ["training", "education", "learning"]):
        return "training_deck"
    if any(kw in all_text for kw in ["organogram", "organization", "hierarchy"]):
        return "reference_doc"
    return "product_brochure"


def extract_file(file_entry):
    """Full extraction of a single file following strict rules"""
    seq = file_entry["sequence"]
    filename = file_entry["filename"]
    file_id = file_entry["file_id"]
    file_type = file_entry["file_type"]
    filepath = os.path.join(BROCHURE_DIR, filename)

    print(f"\n{'='*60}")
    print(f"FILE {seq}: {filename}")
    print(f"  ID: {file_id}")
    print(f"  Type: {file_type}")

    # Rule 2: Confirm source BEFORE extraction
    if not os.path.exists(filepath):
        print(f"  ERROR: File not found at {filepath}")
        return None

    file_size = os.path.getsize(filepath)

    # Extract pages based on file type
    if file_type == "pdf":
        pages = extract_pdf_pages(filepath)
    elif file_type == "pptx":
        pages = extract_pptx_pages(filepath)
    else:
        print(f"  SKIP: Unsupported file type {file_type}")
        return None

    # Classify content
    content_type = classify_content(pages, filename)

    # Build extraction result
    result = {
        "extraction_id": f"{seq:03d}",
        "file_id": file_id,
        "source_file": filename,
        "file_type": file_type,
        "file_fingerprint": {
            "size_bytes": file_size,
            "sha256_12": file_entry.get("sha256_12", ""),
            "processed_at": NOW
        },
        "content_type": content_type,
        "parser_used": "multi_pass",
        "extraction_confidence": 0.0,
        "total_pages": len(pages),
        "page_extractions": [],
        "products": [],
        "key_insights": [],
        "brand_intelligence": {},
        "needs_review": True,
        "extracted_at": NOW
    }

    # Store per-page data (Rule 2: every page tied to source)
    for page in pages:
        page_entry = {
            "source_file": filename,
            "file_id": file_id,
            "page_number": page["page_number"],
            "parser_used": page["extraction_method"],
            "extraction_method": page["extraction_method"],
            "text_length": page["text_length"],
            "is_blank": page["is_blank"],
            "has_tables": page["has_tables"]
        }
        result["page_extractions"].append(page_entry)

    # Store full text for later processing
    result["_raw_text_by_page"] = {
        str(p["page_number"]): p["text"] for p in pages if p["text"]
    }
    result["_tables_by_page"] = {
        str(p["page_number"]): p["tables"] for p in pages if p["tables"]
    }

    # Summary stats
    total_text = sum(p["text_length"] for p in pages)
    blank_pages = sum(1 for p in pages if p["is_blank"])
    table_pages = sum(1 for p in pages if p["has_tables"])

    print(f"  Pages: {len(pages)} ({blank_pages} blank, {table_pages} with tables)")
    print(f"  Total text: {total_text} chars")
    print(f"  Content type: {content_type}")
    print(f"  Parsers used: {set(p['extraction_method'] for p in pages)}")

    result["extraction_confidence"] = 0.9 if total_text > 200 else (0.5 if total_text > 50 else 0.2)

    return result


def process_batch(start_seq, end_seq):
    """Process a batch of files"""
    file_index = load_file_index()
    files = file_index["files"]

    results = []
    for f in files:
        seq = f["sequence"]
        if seq < start_seq or seq > end_seq:
            continue
        if f.get("processed"):
            print(f"  SKIP: File {seq} already processed")
            continue

        result = extract_file(f)
        if result:
            # Save raw extraction
            safe_name = f["filename"].lower()
            safe_name = safe_name.replace(" ", "_").replace(".", "_").replace("-", "_")
            safe_name = safe_name.rsplit("_", 1)[0] if "_" in safe_name else safe_name
            safe_name = safe_name[:40]
            out_path = os.path.join(RAW_DIR, f"{seq:03d}_{safe_name}.json")
            with open(out_path, "w") as fout:
                json.dump(result, fout, indent=2)
            print(f"  Saved: {out_path}")
            results.append(result)

    return results


if __name__ == "__main__":
    import sys
    start = int(sys.argv[1]) if len(sys.argv) > 1 else 26
    end = int(sys.argv[2]) if len(sys.argv) > 2 else 50
    print(f"Processing Files {start}-{end}...")
    results = process_batch(start, end)
    print(f"\nBatch complete: {len(results)} files processed")
